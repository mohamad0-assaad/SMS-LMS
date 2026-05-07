from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.util as util

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Color palette
DARK_BG    = RGBColor(0x0F, 0x17, 0x2A)   # dark navy
ACCENT     = RGBColor(0x6C, 0x63, 0xFF)   # purple-blue
ACCENT2    = RGBColor(0x00, 0xD4, 0xAA)   # teal
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xD6, 0xF1)
YELLOW     = RGBColor(0xFF, 0xD7, 0x00)
MID_BLUE   = RGBColor(0x1E, 0x2D, 0x4E)
CARD_BG    = RGBColor(0x16, 0x23, 0x3E)

def set_bg(slide, color=DARK_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape

def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox

def add_para(tf, text, font_size=16, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, space_before=0, bullet_char=""):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = (bullet_char + " " if bullet_char else "") + text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return p

def accent_bar(slide, y=0.72, width=13.33):
    add_rect(slide, 0, y, width, 0.07, ACCENT)

def header_band(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.3, MID_BLUE)
    accent_bar(slide, 1.23, 13.33)
    add_text_box(slide, title, 0.4, 0.18, 12.5, 0.7,
                 font_size=32, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, subtitle, 0.4, 0.82, 12.5, 0.38,
                     font_size=15, color=ACCENT2, align=PP_ALIGN.LEFT)

# ─────────────────────────────────────────────
# TITLE SLIDE
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
# Gradient-like bands
add_rect(slide, 0, 0, 13.33, 2.5, MID_BLUE)
add_rect(slide, 0, 2.43, 13.33, 0.14, ACCENT)
add_rect(slide, 0, 2.57, 13.33, 0.14, ACCENT2)

# Logo area placeholder — colored square
add_rect(slide, 0.5, 0.35, 1.2, 1.2, ACCENT)
add_text_box(slide, "AI", 0.5, 0.55, 1.2, 0.8, font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text_box(slide, "SkillUp LMS", 2.0, 0.3, 10.8, 1.0, font_size=44, bold=True, color=WHITE)
add_text_box(slide, "AI-Powered School Management & Learning System",
             2.0, 1.1, 10.8, 0.7, font_size=20, color=ACCENT2)

# Info cards
info = [
    ("Presented by:", "Saad Obeid"),
    ("Course:", "Capstone Project"),
    ("Supervisor:", "TBD"),
    ("Date:", "May 2026"),
]
for i, (label, val) in enumerate(info):
    x = 0.5 + i * 3.2
    add_rect(slide, x, 3.1, 2.9, 1.0, CARD_BG)
    add_rect(slide, x, 3.1, 2.9, 0.07, ACCENT)
    add_text_box(slide, label, x+0.1, 3.15, 2.7, 0.35, font_size=11, color=ACCENT2, bold=True)
    add_text_box(slide, val,   x+0.1, 3.48, 2.7, 0.55, font_size=15, color=WHITE, bold=True)

# Bottom tagline
add_text_box(slide, "Full-Stack  ·  TypeScript / React / Node.js  ·  MongoDB  ·  Gemini AI  ·  Python ML",
             0, 6.9, 13.33, 0.5, font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER, italic=True)

# ─────────────────────────────────────────────
# SLIDE 1 — Table of Contents
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
header_band(slide, "Table of Contents")

topics = [
    ("01", "Introduction"),
    ("02", "Problem Statement"),
    ("03", "Solution Overview"),
    ("04", "Technology Stack"),
    ("05", "System Design"),
    ("06", "Implementation Details"),
    ("07", "Challenges Encountered"),
    ("08", "Capstone"),
    ("09", "Future Work"),
    ("10", "Conclusion"),
]
cols = [topics[:5], topics[5:]]
for ci, col in enumerate(cols):
    for ri, (num, topic) in enumerate(col):
        x = 0.5 + ci * 6.5
        y = 1.6 + ri * 1.0
        add_rect(slide, x, y, 6.0, 0.75, CARD_BG)
        add_rect(slide, x, y, 0.55, 0.75, ACCENT)
        add_text_box(slide, num, x+0.02, y+0.1, 0.5, 0.55, font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, topic, x+0.65, y+0.12, 5.2, 0.55, font_size=17, color=WHITE, bold=True)

# ─────────────────────────────────────────────
# SLIDE 2 — Introduction
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
header_band(slide, "Introduction", "What is SkillUp LMS?")

bullets = [
    ("Overview", "SkillUp LMS is a comprehensive full-stack School Management & Learning Management System built for modern educational institutions."),
    ("Motivation", "Traditional school systems lack AI-driven personalization and unified management. SkillUp unifies admin, academic, and AI features in a single platform."),
    ("Objectives", "• Automate school operations (timetables, fees, attendance)\n• Enable AI-driven learning insights for students\n• Provide real-time analytics for teachers and administrators\n• Support four user roles: Admin, Teacher, Student, Parent"),
]
for i, (title, body) in enumerate(bullets):
    y = 1.55 + i * 1.7
    add_rect(slide, 0.4, y, 12.4, 1.5, CARD_BG)
    add_rect(slide, 0.4, y, 0.08, 1.5, ACCENT)
    add_text_box(slide, title, 0.6, y+0.08, 3.0, 0.42, font_size=14, bold=True, color=ACCENT2)
    add_text_box(slide, body,  0.6, y+0.45, 11.8, 0.95, font_size=14, color=LIGHT_GRAY, wrap=True)

# ─────────────────────────────────────────────
# SLIDE 3 — Problem Statement
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
header_band(slide, "Problem Statement", "What problems does SkillUp solve?")

problems = [
    ("Fragmented Tools", "Schools use disconnected software for grades, attendance, fees, and communication — creating data silos and manual overhead."),
    ("No AI Personalization", "Students receive one-size-fits-all education without data-driven feedback on skill gaps, study strategies, or failure risk."),
    ("Limited Parent Visibility", "Parents lack real-time access to their child's academic performance, attendance, and fee status."),
    ("Manual Scheduling", "Timetable generation is time-consuming and error-prone when done manually across subjects, rooms, and teachers."),
]
for i, (title, body) in enumerate(problems):
    col = i % 2
    row = i // 2
    x = 0.4 + col * 6.45
    y = 1.55 + row * 2.3
    add_rect(slide, x, y, 6.1, 2.1, CARD_BG)
    add_rect(slide, x, y, 6.1, 0.08, ACCENT if col == 0 else ACCENT2)
    add_text_box(slide, f"⚠  {title}", x+0.2, y+0.15, 5.7, 0.5, font_size=15, bold=True, color=YELLOW)
    add_text_box(slide, body, x+0.2, y+0.6, 5.7, 1.35, font_size=13, color=LIGHT_GRAY, wrap=True)

# ─────────────────────────────────────────────
# SLIDE 4 — Solution Overview
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
header_band(slide, "Solution Overview", "A unified, AI-powered educational platform")

features = [
    ("Role-Based Dashboards",  "Tailored UX for Admin, Teacher, Student, and Parent with relevant data and actions per role."),
    ("AI-Powered Insights",    "Gemini 2.5 Flash generates exams, class insights, and personalized recommendations."),
    ("ML Skill Predictions",   "Pre-trained scikit-learn models predict student failure risk and recommend study strategies."),
    ("Automated Timetabling",  "Inngest + Gemini AI generates optimized class schedules asynchronously."),
    ("Full Academic Suite",    "Exams, assignments, attendance, report cards, study materials, and messaging — all in one app."),
    ("Financial Management",   "Track student fees, teacher salaries, and school expenses with full payment history."),
]
for i, (title, body) in enumerate(features):
    col = i % 3
    row = i // 3
    x = 0.4 + col * 4.3
    y = 1.55 + row * 2.35
    add_rect(slide, x, y, 4.0, 2.1, CARD_BG)
    add_rect(slide, x, y, 4.0, 0.08, ACCENT)
    add_text_box(slide, title, x+0.15, y+0.15, 3.7, 0.5, font_size=14, bold=True, color=ACCENT2)
    add_text_box(slide, body, x+0.15, y+0.6, 3.7, 1.35, font_size=12, color=LIGHT_GRAY, wrap=True)

# ─────────────────────────────────────────────
# SLIDE 5 — Technology Stack
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
header_band(slide, "Technology Stack", "Languages · Frameworks · Databases · APIs · Tools")

stack = [
    ("Frontend",    ["React 19", "Vite 6", "TypeScript 5.9", "TailwindCSS 4", "React Router 7", "Lucide Icons"]),
    ("Backend",     ["Node.js + Express 5", "TypeScript 5", "JWT Auth", "bcryptjs", "Multer", "Helmet / CORS"]),
    ("Database",    ["MongoDB 7 (Atlas)", "Mongoose 9.3", "Connection Pooling", "16 Collections", "Indexed Queries"]),
    ("AI / ML",     ["Google Gemini 2.5 Flash", "Python FastAPI", "scikit-learn 1.5", "joblib", "NumPy"]),
    ("DevOps/Infra",["Inngest (Event Queue)", "Morgan (logging)", "Vite Dev Server", "REST API", "Git / GitHub"]),
]
for i, (layer, techs) in enumerate(stack):
    x = 0.35 + (i % 3) * 4.25
    y = 1.55 + (i // 3) * 2.55
    w = 3.9
    add_rect(slide, x, y, w, 2.3, CARD_BG)
    add_rect(slide, x, y, w, 0.45, ACCENT if i % 2 == 0 else MID_BLUE)
    add_text_box(slide, layer, x+0.15, y+0.05, w-0.3, 0.38, font_size=14, bold=True, color=WHITE)
    for j, tech in enumerate(techs):
        add_text_box(slide, f"• {tech}", x+0.15, y+0.52+j*0.3, w-0.25, 0.32, font_size=11.5, color=LIGHT_GRAY)

# ─────────────────────────────────────────────
# SLIDE 6 — System Design
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
header_band(slide, "System Design", "Multi-tier architecture with AI microservices")

# Left: Architecture tiers
tiers = [
    (ACCENT,  "Tier 1 — Frontend (React SPA)",          "Port 5175 · Vite · Role-based routing · TailwindCSS"),
    (ACCENT2, "Tier 2 — Backend API (Node/Express)",     "Port 5000 · REST API · JWT auth · 17 route modules"),
    (YELLOW,  "Tier 3 — Database (MongoDB Atlas)",       "16 collections · Mongoose ODM · Cloud / local"),
    (RGBColor(0xFF,0x6B,0x6B), "Tier 4 — AI Services",  "Gemini 2.5 Flash + Python FastAPI ML (port 8000)"),
]
for i, (color, title, desc) in enumerate(tiers):
    y = 1.55 + i * 1.35
    add_rect(slide, 0.4, y, 7.2, 1.18, CARD_BG)
    add_rect(slide, 0.4, y, 0.12, 1.18, color)
    add_text_box(slide, title, 0.65, y+0.08, 6.8, 0.42, font_size=14, bold=True, color=WHITE)
    add_text_box(slide, desc,  0.65, y+0.5,  6.8, 0.55, font_size=12, color=LIGHT_GRAY)
    if i < 3:
        add_text_box(slide, "↓", 3.7, y+1.18, 0.5, 0.22, font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# Right: Data flow
add_rect(slide, 8.1, 1.55, 4.85, 5.42, CARD_BG)
add_rect(slide, 8.1, 1.55, 4.85, 0.08, ACCENT2)
add_text_box(slide, "Data Flow", 8.25, 1.62, 4.6, 0.45, font_size=14, bold=True, color=ACCENT2)

flow_lines = [
    "User Action (Browser)",
    "  ↓  HTTP request",
    "Express Router",
    "  ↓  Auth middleware (JWT)",
    "Controller",
    "  ↓  Query",
    "MongoDB Atlas",
    "",
    "Long tasks → Inngest queue",
    "  → Gemini AI (async)",
    "",
    "Student ML → FastAPI",
    "  → scikit-learn model",
    "  ← failure_risk + strategy",
]
for i, line in enumerate(flow_lines):
    color = ACCENT2 if "↓" in line or "→" in line or "←" in line else (YELLOW if "Inngest" in line or "FastAPI" in line else LIGHT_GRAY)
    bold = line in ("User Action (Browser)", "Express Router", "Controller", "MongoDB Atlas")
    add_text_box(slide, line, 8.25, 2.15+i*0.28, 4.5, 0.3, font_size=11, color=color, bold=bold)

# ─────────────────────────────────────────────
# SLIDE 7 (was 6 in outline) — Implementation Details
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
header_band(slide, "Implementation Details", "Key modules and how they were built")

details = [
    ("AI Exam Generation",
     "Teacher submits topic → Express controller calls Gemini 2.5 Flash → structured MCQ / short-answer JSON returned → saved to MongoDB Exam model → students can attempt exam."),
    ("SkillUp ML Predictions",
     "Student performance data (14 features: scores, attendance, study time, topic mastery) → POST /predict to FastAPI → scikit-learn models return failure_risk, recommended_topic, study_strategy → displayed on StudentSkillInsightPage."),
    ("Timetable Generation",
     "Admin triggers generation → Inngest event queued → background function fetches classes, teachers, subjects → Gemini AI produces conflict-free schedule JSON → stored in Timetable collection."),
    ("Role-Based Auth",
     "JWT issued on login, stored in HttpOnly cookie → auth middleware decodes token and attaches user → role guards (admin/teacher/student/parent) protect each route group."),
]
for i, (title, body) in enumerate(details):
    row = i // 2
    col = i % 2
    x = 0.4 + col * 6.45
    y = 1.55 + row * 2.45
    add_rect(slide, x, y, 6.1, 2.25, CARD_BG)
    add_rect(slide, x, y, 0.1, 2.25, ACCENT if col == 0 else ACCENT2)
    add_text_box(slide, title, x+0.25, y+0.1, 5.7, 0.45, font_size=14, bold=True, color=ACCENT2 if col == 0 else ACCENT)
    add_text_box(slide, body,  x+0.25, y+0.6, 5.7, 1.5,  font_size=12, color=LIGHT_GRAY, wrap=True)

# ─────────────────────────────────────────────
# SLIDE 8 — Challenges Encountered
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
header_band(slide, "Challenges Encountered", "Technical & non-technical hurdles")

challenges = [
    ("Technical", [
        ("Gemini API Rate Limits", "Handled with user-facing error messages and retry guidance when quota is exceeded."),
        ("ML Model Integration",   "Bridging Node.js backend to Python FastAPI required careful CORS and JSON contract management."),
        ("Async Timetabling",      "Inngest event queue needed custom retry logic to handle Gemini timeouts gracefully."),
        ("MongoDB Schema Design",  "16 inter-related collections required careful ref/populate decisions to avoid over-fetching."),
    ]),
    ("Non-Technical", [
        ("Scope Management",   "Balancing 37 frontend pages within deadlines required strict feature prioritization."),
        ("Role Testing",       "Testing 4 distinct user roles meant maintaining multiple test accounts and seed data."),
        ("AI Prompt Tuning",   "Iterating on Gemini prompts to reliably output structured JSON for exams and timetables."),
        ("UI/UX Consistency",  "Ensuring TailwindCSS components remained consistent across all pages and screen sizes."),
    ]),
]
for ci, (category, items) in enumerate(challenges):
    x = 0.4 + ci * 6.45
    add_rect(slide, x, 1.55, 6.1, 0.5, ACCENT if ci == 0 else ACCENT2)
    add_text_box(slide, category + " Challenges", x+0.2, 1.6, 5.7, 0.42, font_size=15, bold=True, color=WHITE)
    for ri, (title, desc) in enumerate(items):
        y = 2.18 + ri * 1.22
        add_rect(slide, x, y, 6.1, 1.1, CARD_BG)
        add_rect(slide, x, y, 6.1, 0.06, ACCENT if ci == 0 else ACCENT2)
        add_text_box(slide, title, x+0.15, y+0.1, 5.8, 0.38, font_size=13, bold=True, color=YELLOW)
        add_text_box(slide, desc,  x+0.15, y+0.5, 5.8, 0.55, font_size=12, color=LIGHT_GRAY, wrap=True)

# ─────────────────────────────────────────────
# SLIDE 9 — Capstone
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
header_band(slide, "Capstone", "Competencies · Course Integration · Real-World Usage")

sections = [
    ("Core Competencies Demonstrated", ACCENT, [
        "Full-stack web development (React + Node.js + MongoDB)",
        "RESTful API design with JWT authentication & RBAC",
        "AI/ML integration (Gemini API + scikit-learn models)",
        "Asynchronous event-driven architecture (Inngest)",
        "UI/UX design with component-based architecture",
        "Database modeling (16 MongoDB collections)",
    ]),
    ("Courses Applied", ACCENT2, [
        "Software Engineering — system design & architecture",
        "Database Management — MongoDB schema & queries",
        "Algorithms & Data Structures — ML model pipelines",
        "Web Development — React, REST APIs, TypeScript",
        "AI / Machine Learning — model training & prediction",
        "Project Management — sprint planning & delivery",
    ]),
    ("Real-World Usage", YELLOW, [
        "Schools can deploy for student enrollment & management",
        "Teachers generate AI exams in seconds, not hours",
        "Students receive personalized study recommendations",
        "Parents track child's attendance, grades & fees online",
        "Admins automate timetabling & financial reporting",
    ]),
]
for ci, (title, color, items) in enumerate(sections):
    x = 0.35 + ci * 4.3
    add_rect(slide, x, 1.55, 4.0, 0.5, color)
    add_text_box(slide, title, x+0.1, 1.59, 3.8, 0.44, font_size=12.5, bold=True, color=DARK_BG)
    add_rect(slide, x, 2.05, 4.0, 5.05, CARD_BG)
    for ri, item in enumerate(items):
        add_text_box(slide, f"✓  {item}", x+0.15, 2.15+ri*0.75, 3.7, 0.65, font_size=12, color=LIGHT_GRAY, wrap=True)

# ─────────────────────────────────────────────
# SLIDE 10 — Future Work
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
header_band(slide, "Future Work", "Limitations, enhancements & scalability roadmap")

future = [
    ("Current Limitations", YELLOW, [
        "ML models trained on synthetic data — need real student datasets",
        "No mobile app (responsive web only)",
        "Gemini API dependency — subject to quota and cost",
        "No real-time notifications (polling-based updates)",
    ]),
    ("Planned Features", ACCENT2, [
        "React Native mobile app for students and parents",
        "Real-time websocket notifications (Socket.io)",
        "Multi-language / localization support (i18n)",
        "Custom ML model training on real institutional data",
        "Advanced analytics dashboard with charts & exports",
        "Integration with Google Classroom / Microsoft Teams",
    ]),
    ("Scalability", ACCENT, [
        "Containerize services with Docker + Docker Compose",
        "Deploy on Kubernetes for horizontal scaling",
        "Migrate to microservices per domain (auth, exams, AI)",
        "CDN for static assets and study materials",
        "GraphQL API layer for flexible client queries",
    ]),
]
for ci, (title, color, items) in enumerate(future):
    x = 0.35 + ci * 4.3
    add_rect(slide, x, 1.55, 4.0, 0.5, color)
    add_text_box(slide, title, x+0.1, 1.59, 3.8, 0.44, font_size=13, bold=True, color=DARK_BG)
    add_rect(slide, x, 2.05, 4.0, 5.0, CARD_BG)
    add_rect(slide, x, 2.05, 0.08, 5.0, color)
    for ri, item in enumerate(items):
        add_text_box(slide, f"→  {item}", x+0.2, 2.15+ri*0.75, 3.6, 0.65, font_size=12, color=LIGHT_GRAY, wrap=True)

# ─────────────────────────────────────────────
# SLIDE 11 — Conclusion
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
header_band(slide, "Conclusion", "Summary · Takeaways · Reflections")

# Impact summary box
add_rect(slide, 0.4, 1.55, 12.4, 1.2, MID_BLUE)
add_rect(slide, 0.4, 1.55, 0.12, 1.2, ACCENT)
add_text_box(slide, "Project Impact",
             0.65, 1.6, 12.0, 0.42, font_size=13, bold=True, color=ACCENT2)
add_text_box(slide,
             "SkillUp LMS delivers a production-ready, AI-augmented school management platform that reduces administrative burden, "
             "personalizes student learning, and connects all stakeholders — admins, teachers, students, and parents — in a single unified system.",
             0.65, 1.95, 12.0, 0.7, font_size=13, color=LIGHT_GRAY, wrap=True)

takeaways = [
    ("AI Accelerates Education",   "Integrating Gemini and ML models transforms static school data into actionable, personalized insights."),
    ("Full-Stack Complexity",       "Building a real-world LMS required coordinating 37 frontend pages, 17 API routes, and 16 database models simultaneously."),
    ("Async is Essential",          "Event-driven patterns (Inngest) were critical for long-running AI operations without blocking the user experience."),
    ("Design for Real Users",       "Role-based UX ensures each user type (admin, teacher, student, parent) sees only what is relevant and actionable for them."),
]
for i, (title, body) in enumerate(takeaways):
    col = i % 2
    row = i // 2
    x = 0.4 + col * 6.45
    y = 3.05 + row * 1.7
    add_rect(slide, x, y, 6.1, 1.5, CARD_BG)
    add_rect(slide, x, y, 6.1, 0.07, ACCENT if col == 0 else ACCENT2)
    add_text_box(slide, title, x+0.2, y+0.13, 5.7, 0.42, font_size=13, bold=True, color=YELLOW)
    add_text_box(slide, body,  x+0.2, y+0.58, 5.7, 0.82, font_size=12, color=LIGHT_GRAY, wrap=True)

# Final remark
add_rect(slide, 0.4, 6.42, 12.4, 0.78, CARD_BG)
add_rect(slide, 0.4, 6.42, 12.4, 0.07, ACCENT2)
add_text_box(slide,
             '"SkillUp LMS proves that modern web technologies and AI can be combined to build scalable, impactful tools for real educational challenges."',
             0.6, 6.52, 12.1, 0.6, font_size=13, italic=True, color=ACCENT2, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# Save
# ─────────────────────────────────────────────
out = r"c:\Users\user\New folder\SMS-LMS\SkillUp_LMS_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
